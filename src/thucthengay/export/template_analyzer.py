"""PowerPoint template shape inventory and placeholder matching."""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE

from thucthengay.models import (
    PlaceholderType,
    TemplatePlaceholder,
    TemplatePlaceholderSelector,
)


@dataclass(frozen=True)
class TemplateShape:
    """Stable subset of PPTX shape metadata used by placeholder resolution."""

    element_id: int
    name: str
    title: str | None
    descr: str | None
    text: str
    shape_type: str
    has_text_frame: bool
    is_picture: bool
    left: int
    top: int
    width: int
    height: int

    @property
    def area(self) -> int:
        return self.width * self.height


@dataclass(frozen=True)
class PlaceholderMatch:
    """Resolved placeholder and the evidence used for the match."""

    placeholder: TemplatePlaceholder
    shape: TemplateShape | None
    method: str
    configured_element_id: int | None
    candidates: tuple[TemplateShape, ...] = ()

    @property
    def resolved(self) -> bool:
        return self.shape is not None

    @property
    def ambiguous(self) -> bool:
        return len(self.candidates) > 1


def build_shape_inventory(shapes: Any) -> tuple[TemplateShape, ...]:
    """Return a recursive inventory of all slide shapes."""
    inventory: list[TemplateShape] = []
    for shape in shapes:
        inventory.append(_shape_metadata(shape))
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            inventory.extend(build_shape_inventory(child_shapes))
    return tuple(inventory)


def resolve_placeholder(
    placeholder: TemplatePlaceholder,
    inventory: tuple[TemplateShape, ...],
) -> PlaceholderMatch:
    """Resolve one configured placeholder against the current PPTX inventory."""
    if placeholder.selector is not None:
        candidates = _selector_candidates(placeholder.selector, inventory)
        match = _single_match(placeholder, candidates, method="selector")
        if match is not None:
            return match

    conventional = TemplatePlaceholderSelector(name=f"ttn:{placeholder.field}")
    candidates = _selector_candidates(conventional, inventory)
    match = _single_match(placeholder, candidates, method="field_name")
    if match is not None:
        return match

    if placeholder.kind == PlaceholderType.TEXT:
        candidates = tuple(
            shape
            for shape in inventory
            if shape.has_text_frame and _normalized(shape.text) == _normalized(placeholder.field)
        )
        match = _single_match(placeholder, candidates, method="field_text")
        if match is not None:
            return match

    if placeholder.element_id is not None:
        shape = _shape_by_id(inventory).get(placeholder.element_id)
        if shape is not None:
            return PlaceholderMatch(
                placeholder=_placeholder_for_shape(placeholder, shape),
                shape=shape,
                method="configured_element_id",
                configured_element_id=placeholder.element_id,
            )

    return PlaceholderMatch(
        placeholder=placeholder,
        shape=None,
        method="missing",
        configured_element_id=placeholder.element_id,
    )


def shape_inventory_payload(inventory: tuple[TemplateShape, ...]) -> list[dict[str, Any]]:
    """Serialize inventory into template metadata diagnostics."""
    return [
        {
            "id": str(shape.element_id),
            "name": shape.name,
            "title": shape.title,
            "descr": shape.descr,
            "text": shape.text,
            "shape_type": shape.shape_type,
            "has_text_frame": shape.has_text_frame,
            "is_picture": shape.is_picture,
            "bbox": {
                "emu": {
                    "x": shape.left,
                    "y": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
            },
        }
        for shape in inventory
    ]


def _shape_metadata(shape: Any) -> TemplateShape:
    xml = _c_nv_pr(shape)
    shape_type = getattr(shape, "shape_type", None)
    return TemplateShape(
        element_id=int(shape.shape_id),
        name=getattr(shape, "name", "") or "",
        title=xml.get("title") if xml is not None else None,
        descr=xml.get("descr") if xml is not None else None,
        text=_shape_text(shape),
        shape_type=_shape_type_name(shape_type),
        has_text_frame=bool(getattr(shape, "has_text_frame", False)),
        is_picture=shape_type == MSO_SHAPE_TYPE.PICTURE,
        left=int(shape.left),
        top=int(shape.top),
        width=int(shape.width),
        height=int(shape.height),
    )


def _c_nv_pr(shape: Any) -> Any | None:
    nodes = shape.element.xpath("./*/p:cNvPr")
    if not nodes:
        nodes = shape.element.xpath(".//p:cNvPr")
    return nodes[0] if nodes else None


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def _shape_type_name(shape_type: Any) -> str:
    try:
        return MSO_SHAPE_TYPE(shape_type).name
    except Exception:
        return str(shape_type)


def _selector_candidates(
    selector: TemplatePlaceholderSelector,
    inventory: tuple[TemplateShape, ...],
) -> tuple[TemplateShape, ...]:
    return tuple(shape for shape in inventory if _selector_matches(selector, shape))


def _selector_matches(selector: TemplatePlaceholderSelector, shape: TemplateShape) -> bool:
    if selector.name is not None and shape.name != selector.name:
        return False
    if selector.title is not None and shape.title != selector.title:
        return False
    if selector.descr is not None and shape.descr != selector.descr:
        return False
    if selector.alt_text is not None and shape.descr != selector.alt_text:
        return False
    if selector.text is not None and _normalized(shape.text) != _normalized(selector.text):
        return False
    return True


def _single_match(
    placeholder: TemplatePlaceholder,
    candidates: tuple[TemplateShape, ...],
    *,
    method: str,
) -> PlaceholderMatch | None:
    if len(candidates) == 1:
        shape = candidates[0]
        return PlaceholderMatch(
            placeholder=_placeholder_for_shape(placeholder, shape),
            shape=shape,
            method=method,
            configured_element_id=placeholder.element_id,
        )
    if len(candidates) > 1:
        return PlaceholderMatch(
            placeholder=placeholder,
            shape=None,
            method=method,
            configured_element_id=placeholder.element_id,
            candidates=candidates,
        )
    return None


def _placeholder_for_shape(
    placeholder: TemplatePlaceholder,
    shape: TemplateShape,
) -> TemplatePlaceholder:
    return placeholder.model_copy(
        update={
            "element_id": shape.element_id,
            "diagnostic_name": placeholder.diagnostic_name or shape.name,
        }
    )


def _shape_by_id(inventory: tuple[TemplateShape, ...]) -> dict[int, TemplateShape]:
    return {shape.element_id: shape for shape in inventory}


def _normalized(value: str) -> str:
    return " ".join(value.casefold().split())
