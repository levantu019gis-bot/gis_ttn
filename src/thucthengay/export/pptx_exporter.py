"""Headless combined PowerPoint export service."""

from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path

from pptx import Presentation
from pydantic import ValidationError

from thucthengay.export.compare_text import resolve_compare_text_panes
from thucthengay.export.pptx_slide_copy import (
    copy_only_slide,
    find_shape_by_element_id,
    replace_shape_with_picture,
    replace_text,
)
from thucthengay.export.preflight import build_export_preflight_plan
from thucthengay.export.txt_values import (
    SUPPORTED_TEXT_FIELDS,
    TxtLineResolution,
    resolve_export_text,
)
from thucthengay.models import (
    Composition,
    ExportedComposition,
    ExportPptxResult,
    ExportPptxSummary,
    Issue,
    IssueScope,
    IssueSeverity,
    PlaceholderType,
    TargetConfig,
    TemplateMetadata,
    TemplatePlaceholder,
)
from thucthengay.workspace import WorkspaceService


def export_combined_pptx(
    workspace_service: WorkspaceService,
    targets: Iterable[TargetConfig],
    *,
    output_path: str | Path,
    template_issues: Iterable[Issue] = (),
    composition_ids: Iterable[str] | None = None,
) -> ExportPptxResult:
    """Write one combined PPTX for included compositions sorted by review order."""
    target_map = {target.id: target for target in targets}
    selected_ids = set(composition_ids) if composition_ids is not None else None
    included = [
        composition
        for composition in workspace_service.list_compositions()
        if composition.include
        and (selected_ids is None or composition.composition_id in selected_ids)
    ]
    included.sort(key=_export_sort_key)
    selected_target_ids = {composition.target_id for composition in included}
    plan = build_export_preflight_plan(
        workspace_service,
        target_map.values(),
        template_issues=template_issues,
    )
    issues = _issues_for_selected(plan.issues, selected_ids, selected_target_ids)
    if not included:
        return _blocked([_no_exportable_compositions_issue()])
    placeholder_issues = _placeholder_issues(workspace_service, included, target_map)
    issues.extend(placeholder_issues)
    skipped_composition_ids, global_blocking_issues = _skippable_composition_ids(
        issues,
        included,
    )
    if global_blocking_issues:
        return _blocked(issues)
    exportable = [
        composition
        for composition in included
        if composition.composition_id not in skipped_composition_ids
    ]
    if not exportable:
        return _blocked(issues or [_no_exportable_compositions_issue()])

    try:
        resolved_output = _resolve_output_path(workspace_service, output_path)
        resolved_output.parent.mkdir(parents=True, exist_ok=True)
    except (OSError, ValueError) as error:
        return _blocked([_write_issue(f"Khong ghi duoc file PPTX export: {error}")])
    relative_output = _workspace_relative(workspace_service, resolved_output)

    destination = Presentation()
    exported: list[ExportedComposition] = []
    for slide_number, composition in enumerate(exportable, start=1):
        target = target_map[composition.target_id]
        template = _template_metadata(target)
        if slide_number == 1:
            source_for_size = Presentation(template.template_pptx)
            destination.slide_width = source_for_size.slide_width
            destination.slide_height = source_for_size.slide_height
        source = Presentation(template.template_pptx)
        slide = copy_only_slide(source, destination)
        render_path = _resolved_render_path(workspace_service, composition)
        for placeholder in template.placeholders:
            if placeholder.element_id is None:
                continue
            if placeholder.kind == PlaceholderType.MAP_IMAGE:
                replace_shape_with_picture(slide, placeholder.element_id, render_path)
            elif placeholder.kind == PlaceholderType.TEXT:
                resolution = _placeholder_text_resolution(
                    placeholder,
                    composition,
                    target,
                    slide_number,
                    workspace_service=workspace_service,
                )
                replace_text(
                    slide,
                    placeholder.element_id,
                    resolution.text if resolution.ok else "",
                    fit_shape_to_text=_should_fit_shape_to_text(placeholder),
                )
        exported.append(
            ExportedComposition(
                composition_id=composition.composition_id,
                target_id=composition.target_id,
                slide_number=slide_number,
                render_path=composition.artifacts.final_render_path or "",
            )
        )

    try:
        destination.save(resolved_output)
    except OSError as error:
        try:
            resolved_output.unlink(missing_ok=True)
        except OSError:
            pass
        return _blocked([_write_issue(f"Khong ghi duoc file PPTX export: {error}")])
    summary = ExportPptxSummary(
        slide_count=len(exported),
        target_count=len({item.target_id for item in exported}),
        error_count=sum(1 for issue in issues if issue.severity == IssueSeverity.ERROR),
        warning_count=sum(1 for issue in issues if issue.severity == IssueSeverity.WARNING),
    )
    return ExportPptxResult(
        ok=True,
        pptx_path=relative_output,
        exported=exported,
        issues=issues,
        summary=summary,
    )


def _issues_for_selected(
    issues: Iterable[Issue],
    selected_ids: set[str] | None,
    selected_target_ids: set[str],
) -> list[Issue]:
    if selected_ids is None:
        return list(issues)
    selected_issues: list[Issue] = []
    for issue in issues:
        if issue.composition_id is not None and issue.composition_id not in selected_ids:
            continue
        if issue.composition_id is None and issue.target_id is not None:
            if issue.target_id not in selected_target_ids:
                continue
        selected_issues.append(issue)
    return selected_issues


def _skippable_composition_ids(
    issues: Iterable[Issue],
    compositions: list[Composition],
) -> tuple[set[str], list[Issue]]:
    compositions_by_target: dict[str, list[Composition]] = {}
    for composition in compositions:
        compositions_by_target.setdefault(composition.target_id, []).append(composition)
    composition_ids = {composition.composition_id for composition in compositions}
    skipped_ids: set[str] = set()
    global_blocking: list[Issue] = []
    for issue in issues:
        if not issue.blocking:
            continue
        if issue.composition_id is not None:
            if issue.composition_id in composition_ids:
                skipped_ids.add(issue.composition_id)
            else:
                global_blocking.append(issue)
            continue
        if issue.target_id is not None:
            target_compositions = compositions_by_target.get(issue.target_id, [])
            if target_compositions:
                skipped_ids.update(
                    composition.composition_id for composition in target_compositions
                )
            else:
                global_blocking.append(issue)
            continue
        global_blocking.append(issue)
    return skipped_ids, global_blocking


def _placeholder_issues(
    workspace_service: WorkspaceService,
    compositions: list[Composition],
    target_map: dict[str, TargetConfig],
) -> list[Issue]:
    issues: list[Issue] = []
    for slide_number, composition in enumerate(compositions, start=1):
        target = target_map.get(composition.target_id)
        if target is None:
            issues.append(
                _issue(
                    "export.target_missing",
                    "Khong tim thay target cua composition khi export PPTX.",
                    "Kiem tra lai config target va workspace composition.",
                    composition=composition,
                )
            )
            continue
        try:
            template = _template_metadata(target)
        except ValueError as error:
            issues.append(
                _issue(
                    "export.pptx_template_metadata_invalid",
                    "Khong doc duoc template metadata cua target khi export PPTX.",
                    f"Chay lai load config/template truoc khi export. Chi tiet: {error}",
                    composition=composition,
                )
            )
            continue
        render_path = _try_resolved_render_path(composition, target, target_map, slide_number)
        if render_path is not None:
            issues.append(render_path)
        issues.extend(
            _template_placeholder_issues(
                composition,
                target,
                template,
                slide_number,
                workspace_service=workspace_service,
            )
        )
    return issues


def _template_placeholder_issues(
    composition: Composition,
    target: TargetConfig,
    template: TemplateMetadata,
    slide_number: int,
    *,
    workspace_service: WorkspaceService,
) -> list[Issue]:
    source = Presentation(template.template_pptx)
    if len(source.slides) != 1:
        return [
            _issue(
                "export.pptx_template_slide_count_invalid",
                "PPTX template khong co dung mot slide khi export.",
                "Sua template target thanh file PPTX mot slide roi chay lai preflight.",
                composition=composition,
            )
        ]
    issues: list[Issue] = []
    for placeholder in template.placeholders:
        if placeholder.element_id is None:
            if placeholder.required:
                issues.append(
                    _issue(
                        "export.pptx_placeholder_element_missing",
                        (
                            "PPTX placeholder bat buoc chua duoc resolve thanh element id "
                            f"cho field `{placeholder.field}`."
                        ),
                        (
                            "Tai lai config de resolver doc PPTX template, hoac bo sung "
                            "`selector`/`element_id` hop le cho placeholder."
                        ),
                        composition=composition,
                    )
                )
            continue
        source_shape = find_shape_by_element_id(source.slides[0], placeholder.element_id)
        if source_shape is None:
            if placeholder.required:
                issues.append(
                    _issue(
                        "export.pptx_placeholder_element_missing",
                        (
                            "PPTX template thieu element id bat buoc "
                            f"`{placeholder.element_id}` cho field `{placeholder.field}`."
                        ),
                        (
                            "Cap nhat element_id trong target export config theo dung shape id "
                            "cua template PPTX."
                        ),
                        composition=composition,
                    )
                )
            continue
        if placeholder.kind == PlaceholderType.TEXT:
            if not getattr(source_shape, "has_text_frame", False):
                if placeholder.required:
                    issues.append(
                        _issue(
                            "export.pptx_text_placeholder_invalid",
                            (
                                "PPTX element id bat buoc cho text placeholder khong phai "
                                "text shape."
                            ),
                            "Doi placeholder sang text box hoac cap nhat element_id dung.",
                            composition=composition,
                        )
                    )
                continue
            resolution = _placeholder_text_resolution(
                placeholder,
                composition,
                target,
                slide_number,
                workspace_service=workspace_service,
            )
            if resolution.problems:
                if placeholder.required:
                    fields = ", ".join(problem.field for problem in resolution.problems)
                    issues.append(
                        _issue(
                            "export.pptx_placeholder_unresolved",
                            (
                                "Khong resolve duoc text placeholder bat buoc "
                                f"`{fields}` cho PPTX."
                            ),
                            (
                                "Sua field/value placeholder trong target export config hoac bo "
                                "sung du lieu composition/target can thiet."
                            ),
                            composition=composition,
                        )
                    )
    return issues


def _placeholder_text_resolution(
    placeholder: TemplatePlaceholder,
    composition: Composition,
    target: TargetConfig,
    slide_number: int,
    *,
    workspace_service: WorkspaceService,
) -> TxtLineResolution:
    template = placeholder.value if placeholder.value is not None else f"{{{placeholder.field}}}"
    pane_a, pane_b = resolve_compare_text_panes(workspace_service, composition)
    return resolve_export_text(
        template,
        composition,
        target,
        slide_number=slide_number,
        pane_a_composition=pane_a,
        pane_b_composition=pane_b,
        supported_fields=SUPPORTED_TEXT_FIELDS,
        unknown_issue_id="export.pptx_placeholder_unknown",
        unresolved_issue_id="export.pptx_placeholder_unresolved",
    )


def _should_fit_shape_to_text(placeholder: TemplatePlaceholder) -> bool:
    return placeholder.field in {"time", "time_label", "time_label_pane_A", "time_label_pane_B"}


def _template_metadata(target: TargetConfig) -> TemplateMetadata:
    try:
        return TemplateMetadata.model_validate(target.metadata["template_metadata"])
    except (KeyError, ValidationError) as error:
        msg = "target is missing derived template_metadata"
        raise ValueError(msg) from error


def _resolved_render_path(workspace_service: WorkspaceService, composition: Composition) -> Path:
    if not composition.artifacts.final_render_path:
        msg = "composition is missing final_render_path"
        raise ValueError(msg)
    return workspace_service.paths.root / composition.artifacts.final_render_path


def _try_resolved_render_path(
    composition: Composition,
    target: TargetConfig,
    target_map: dict[str, TargetConfig],
    slide_number: int,
) -> Issue | None:
    del target, target_map, slide_number
    if composition.artifacts.final_render_path:
        return None
    return _issue(
        "export.final_render_missing",
        "Composition chua co anh final render de dua vao PPTX.",
        "Chay render final cho composition nay truoc khi export PPTX.",
        composition=composition,
    )


def _resolve_output_path(workspace_service: WorkspaceService, output_path: str | Path) -> Path:
    path = Path(output_path)
    if not path.is_absolute():
        path = workspace_service.paths.root / path
    resolved = path.resolve()
    root = workspace_service.paths.root.resolve()
    try:
        resolved.relative_to(root)
    except ValueError as error:
        msg = f"PPTX output path must stay inside workspace: {output_path!r}"
        raise ValueError(msg) from error
    return resolved


def _workspace_relative(workspace_service: WorkspaceService, path: Path) -> str:
    return path.resolve().relative_to(workspace_service.paths.root.resolve()).as_posix()


def _blocked(issues: list[Issue]) -> ExportPptxResult:
    return ExportPptxResult(
        ok=False,
        issues=issues,
        summary=ExportPptxSummary(
            error_count=sum(1 for issue in issues if issue.severity == IssueSeverity.ERROR),
            warning_count=sum(1 for issue in issues if issue.severity == IssueSeverity.WARNING),
        ),
    )


def _issue(
    issue_id: str,
    message: str,
    remediation: str,
    *,
    composition: Composition,
) -> Issue:
    return Issue(
        issue_id=issue_id,
        severity=IssueSeverity.ERROR,
        scope=IssueScope.EXPORT,
        target_id=composition.target_id,
        composition_id=composition.composition_id,
        message=message,
        remediation=remediation,
    )


def _write_issue(message: str) -> Issue:
    return Issue(
        issue_id="export.pptx_write_failed",
        severity=IssueSeverity.ERROR,
        scope=IssueScope.EXPORT,
        message=message,
        remediation=(
            "Chon duong dan PPTX trong workspace, dong file dang khoa, "
            "kiem tra quyen ghi, roi chay export lai."
        ),
    )


def _no_exportable_compositions_issue() -> Issue:
    return Issue(
        issue_id="export.no_exportable_compositions",
        severity=IssueSeverity.ERROR,
        scope=IssueScope.EXPORT,
        message="Khong co composition nao du dieu kien de xuat PPTX.",
        remediation="Sua cac loi blocking trong Preflight hoac bo include cac composition loi.",
    )


def _export_sort_key(composition: Composition) -> tuple[int, int, str]:
    return (
        1 if composition.review_order is None else 0,
        composition.review_order or 0,
        composition.composition_id,
    )
