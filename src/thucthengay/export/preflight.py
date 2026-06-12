"""Headless export preflight planning."""

from __future__ import annotations

from collections import defaultdict
from collections.abc import Iterable
from datetime import date
from pathlib import Path

from pptx import Presentation
from pydantic import ValidationError

from thucthengay.export.final_render import final_render_currentness_issue
from thucthengay.export.txt_values import resolve_txt_line, time_label
from thucthengay.models import (
    Composition,
    ExportPlanRow,
    ExportPreflightPlan,
    ExportPreflightState,
    ExportPreflightSummary,
    Issue,
    IssueScope,
    IssueSeverity,
    PlaceholderType,
    TargetConfig,
    TemplateMetadata,
)
from thucthengay.validation import ValidationContext, validate_export_preflight
from thucthengay.workspace import WorkspaceService


def build_export_preflight_plan(
    workspace_service: WorkspaceService,
    targets: Iterable[TargetConfig],
    *,
    template_issues: Iterable[Issue] = (),
) -> ExportPreflightPlan:
    """Build a full preflight result for included compositions without UI imports."""
    target_map = {target.id: target for target in targets}
    included = [
        composition
        for composition in workspace_service.list_compositions()
        if composition.include
    ]
    included.sort(key=_export_sort_key)

    context_issues: list[Issue] = []
    contexts = [
        _validation_context_for(composition, target_map.get(composition.target_id))
        for composition in included
    ]
    if contexts:
        context_issues.extend(validate_export_preflight(contexts, template_issues).issues)
    else:
        context_issues.extend(template_issues)
    context_issues.extend(_pptx_template_dimension_issues(included, target_map))

    row_issue_map: dict[str, list[Issue]] = defaultdict(list)
    target_issue_map: dict[str, list[Issue]] = defaultdict(list)
    global_issues: list[Issue] = []
    for issue in context_issues:
        if issue.composition_id:
            row_issue_map[issue.composition_id].append(issue)
        elif issue.target_id:
            target_issue_map[issue.target_id].append(issue)
        else:
            global_issues.append(issue)

    rows: list[ExportPlanRow] = []
    all_issues = list(context_issues)
    for index, composition in enumerate(included, start=1):
        target = target_map.get(composition.target_id)
        row_issues = list(row_issue_map.get(composition.composition_id, ()))
        row_issues.extend(target_issue_map.get(composition.target_id, ()))
        row_issues.extend(global_issues)
        row_issues.extend(_composition_export_issues(composition, workspace_service, target))
        row_issues.extend(_target_export_issues(composition, target))
        row_issues.extend(_txt_template_issues(composition, target, slide_number=index))
        all_issues.extend(issue for issue in row_issues if issue not in all_issues)
        rows.append(
            ExportPlanRow(
                composition_id=composition.composition_id,
                target_id=composition.target_id,
                slide_number=index if composition.review_order is not None else None,
                review_order=composition.review_order,
                target_label=_target_label(target, composition.target_id),
                date_label=composition.capture_date.isoformat(),
                time_label=time_label(composition, target=target) if target is not None else "",
                template_status=_template_status(row_issues),
                final_render_path=composition.artifacts.final_render_path,
                issues=row_issues,
            )
        )

    summary = _summary(rows, all_issues)
    return ExportPreflightPlan(rows=rows, issues=all_issues, summary=summary)


def _validation_context_for(
    composition: Composition,
    target: TargetConfig | None,
) -> ValidationContext:
    template_metadata: TemplateMetadata | None = None
    template_error: str | None = None
    if target is not None and "template_metadata" in target.metadata:
        try:
            template_metadata = TemplateMetadata.model_validate(
                target.metadata["template_metadata"]
            )
        except Exception as error:  # noqa: BLE001
            template_error = str(error)
    return ValidationContext(
        target=target,
        composition=composition,
        template_metadata=template_metadata,
        template_metadata_error=template_error,
        require_current_validation=True,
    )


def _composition_export_issues(
    composition: Composition,
    workspace_service: WorkspaceService,
    target: TargetConfig | None,
) -> list[Issue]:
    issues: list[Issue] = []
    if composition.review_order is None:
        issues.append(
            Issue(
                issue_id="export.review_order_missing",
                severity=IssueSeverity.ERROR,
                scope=IssueScope.EXPORT,
                target_id=composition.target_id,
                composition_id=composition.composition_id,
                message="Composition include chua co review_order.",
                remediation="Quay lai Review/Edit va Include/Validate lai composition nay.",
            )
        )
    if not composition.ready:
        issues.append(
            Issue(
                issue_id="export.composition_not_ready",
                severity=IssueSeverity.ERROR,
                scope=IssueScope.COMPOSITION,
                target_id=composition.target_id,
                composition_id=composition.composition_id,
                message="Composition da include nhung chua o trang thai ready.",
                remediation="Chay Include/Validate lai hoac bo include composition nay.",
            )
        )

    render_issue = final_render_currentness_issue(
        workspace_service=workspace_service,
        composition=composition,
        target=target,
    )
    if render_issue is not None:
        issues.append(render_issue)
    return issues


def _target_export_issues(
    composition: Composition,
    target: TargetConfig | None,
) -> list[Issue]:
    if target is None:
        return []
    required_map = [
        placeholder
        for placeholder in target.export.placeholders
        if placeholder.required and placeholder.kind == PlaceholderType.MAP_IMAGE
    ]
    if required_map:
        return []
    return [
        Issue(
            issue_id="export.map_placeholder_missing",
            severity=IssueSeverity.ERROR,
            scope=IssueScope.TEMPLATE,
            target_id=target.id,
            composition_id=composition.composition_id,
            message="Target chua co placeholder map_image bat buoc cho export.",
            remediation="Bo sung `export.placeholders` voi `kind: map_image` va element_id dung.",
        )
    ]


def _pptx_template_dimension_issues(
    compositions: list[Composition],
    target_map: dict[str, TargetConfig],
) -> list[Issue]:
    target_ids = _included_target_ids(compositions)
    templates: dict[str, TemplateMetadata] = {}
    for target_id in target_ids:
        target = target_map.get(target_id)
        if target is None:
            continue
        try:
            templates[target_id] = TemplateMetadata.model_validate(
                target.metadata["template_metadata"]
            )
        except (KeyError, ValidationError):
            continue
    if len({template.template_pptx for template in templates.values()}) <= 1:
        return []

    dimensions: dict[str, tuple[int, int]] = {}
    paths: dict[str, str] = {}
    for target_id, template in templates.items():
        paths[target_id] = template.template_pptx
        if not Path(template.template_pptx).is_file():
            continue
        try:
            presentation = Presentation(template.template_pptx)
        except Exception:  # noqa: BLE001
            continue
        dimensions[target_id] = (int(presentation.slide_width), int(presentation.slide_height))

    if len(dimensions) <= 1:
        return []

    issues: list[Issue] = []
    base_target_id = next(iter(dimensions))
    base_dimension = dimensions[base_target_id]
    for target_id, dimension in dimensions.items():
        if target_id == base_target_id or dimension == base_dimension:
            continue
        issues.append(
            Issue(
                issue_id="export.pptx_template_slide_size_mismatch",
                severity=IssueSeverity.ERROR,
                scope=IssueScope.TEMPLATE,
                target_id=target_id,
                message=(
                    "PPTX template cua cac target include khac kich thuoc slide: "
                    f"`{base_target_id}` {base_dimension[0]}x{base_dimension[1]} EMU, "
                    f"`{target_id}` {dimension[0]}x{dimension[1]} EMU."
                ),
                remediation=(
                    "Dung cac PPTX template co cung slide size/aspect cho mot lan export tong hop, "
                    "hoac tach thanh cac lan export rieng. "
                    f"Template tham chieu: {paths.get(base_target_id, '')}; "
                    f"template lech: {paths.get(target_id, '')}."
                ),
            )
        )
    return issues


def _included_target_ids(compositions: list[Composition]) -> list[str]:
    target_ids: list[str] = []
    for composition in compositions:
        if composition.target_id not in target_ids:
            target_ids.append(composition.target_id)
    return target_ids


def _txt_template_issues(
    composition: Composition,
    target: TargetConfig | None,
    *,
    slide_number: int,
) -> list[Issue]:
    if target is None:
        return []
    template = target.export.template_txt_value
    if not template:
        return [
            Issue(
                issue_id="export.txt_template_missing",
                severity=IssueSeverity.ERROR,
                scope=IssueScope.EXPORT,
                target_id=target.id,
                composition_id=composition.composition_id,
                message="Target chua cau hinh template_txt_value.",
                remediation="Bo sung `export.template_txt_value` cho target truoc khi export TXT.",
            )
        ]

    issues: list[Issue] = []
    resolution = resolve_txt_line(
        template,
        composition,
        target,
        slide_number=slide_number,
    )
    for problem in resolution.problems:
        if problem.issue_id == "export.txt_placeholder_unknown":
            issues.append(
                Issue(
                    issue_id="export.txt_placeholder_unknown",
                    severity=IssueSeverity.ERROR,
                    scope=IssueScope.EXPORT,
                    target_id=target.id,
                    composition_id=composition.composition_id,
                    message=f"TXT template dung placeholder chua ho tro: {problem.field}.",
                    remediation="Sua template_txt_value de chi dung cac placeholder da ho tro.",
                )
            )
        elif problem.issue_id == "export.txt_time_label_unresolved":
            issues.append(
                Issue(
                    issue_id="export.txt_time_label_unresolved",
                    severity=IssueSeverity.ERROR,
                    scope=IssueScope.EXPORT,
                    target_id=target.id,
                    composition_id=composition.composition_id,
                    message=(
                        "TXT template can time_label nhung khong co layer visible hop le "
                        "co thoi gian."
                    ),
                    remediation=(
                        "Quay lai Review/Edit Metadata de sua capture_time va metadata_status "
                        "cua layer visible, hoac danh dau placeholder la optional bang "
                        "`{time_label?}`."
                    ),
                )
            )
        else:
            issues.append(
                Issue(
                    issue_id="export.txt_placeholder_unresolved",
                    severity=IssueSeverity.ERROR,
                    scope=IssueScope.EXPORT,
                    target_id=target.id,
                    composition_id=composition.composition_id,
                    message=f"TXT template khong resolve duoc placeholder: {problem.field}.",
                    remediation=(
                        "Sua metadata/composition/target config de co gia tri "
                        "truoc khi export."
                    ),
                )
            )
    return issues


def _export_sort_key(composition: Composition) -> tuple[int, int, date, str]:
    return (
        1 if composition.review_order is None else 0,
        composition.review_order or 0,
        composition.capture_date,
        composition.composition_id,
    )


def _target_label(target: TargetConfig | None, fallback_target_id: str) -> str:
    if target is None:
        return fallback_target_id
    return target.title or target.alias or target.name or target.id


def _template_status(issues: Iterable[Issue]) -> str:
    template_issues = [issue for issue in issues if issue.scope == IssueScope.TEMPLATE]
    if any(issue.severity == IssueSeverity.ERROR for issue in template_issues):
        return "ERROR"
    if template_issues:
        return "WARN"
    return "OK"


def _summary(rows: list[ExportPlanRow], issues: list[Issue]) -> ExportPreflightSummary:
    warning_count = sum(1 for issue in issues if issue.severity == IssueSeverity.WARNING)
    error_count = sum(1 for issue in issues if issue.severity == IssueSeverity.ERROR)
    state = ExportPreflightState.READY
    if error_count:
        state = ExportPreflightState.BLOCKED
    elif warning_count:
        state = ExportPreflightState.WARNING
    return ExportPreflightSummary(
        included_slide_count=len(rows),
        target_count=len({row.target_id for row in rows}),
        skipped_count=sum(1 for row in rows if row.blocking),
        warning_count=warning_count,
        error_count=error_count,
        state=state,
    )
