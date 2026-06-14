from __future__ import annotations

import os
import time as time_module
from datetime import date, time
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")

from PySide6.QtCore import Qt
from PySide6.QtWidgets import QApplication

from thucthengay.config import ConfigLoadResult
from thucthengay.editor.app_shell import AppShell
from thucthengay.editor.models.export_plan_model import ExportPlanRole
from thucthengay.editor.modes.export_mode import ExportMode
from thucthengay.editor.preferences import PreferencesService
from thucthengay.export import (
    ExportHistorySyncResult,
    ExportProgress,
    FullExportResult,
    ensure_final_renders_for_export,
    run_full_export,
)
from thucthengay.models import (
    Composition,
    CompositionArtifacts,
    ExportCompletionState,
    ExportCompletionSummary,
    ExportFinalRenderResult,
    ExportLogWriteResult,
    ExportPptxResult,
    ExportTxtResult,
    GridConfig,
    GridInterval,
    ImageLayer,
    Issue,
    IssueScope,
    IssueSeverity,
    MapFrame,
    MetadataStatus,
    PlaceholderType,
    ProjectConfig,
    TargetConfig,
    TemplateMetadata,
    TemplatePlaceholder,
    ViewState,
)
from thucthengay.render import RasterRenderResult, RenderSpec
from thucthengay.workspace import WorkspaceService


def qapp() -> QApplication:
    return QApplication.instance() or QApplication([])


def target_config(*, txt_template: str = "{slide_number}|{target_id}|{time_label}") -> TargetConfig:
    return TargetConfig(
        id="alpha",
        name="Alpha Target",
        geojson_file="targets/alpha.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": "templates/alpha.pptx",
            "txt_line_template": txt_template,
            "placeholders": [
                {
                    "field": "map",
                    "element_id": 10,
                    "kind": PlaceholderType.MAP_IMAGE,
                    "required": True,
                }
            ],
        },
        metadata={
            "template_metadata": TemplateMetadata(
                template_pptx="templates/alpha.pptx",
                slide_index=0,
                map_frame=MapFrame(x=0, y=0, width=640, height=360),
                placeholders=[
                    TemplatePlaceholder(
                        field="map",
                        element_id=10,
                        kind=PlaceholderType.MAP_IMAGE,
                        required=True,
                    )
                ],
            ).model_dump(mode="json")
        },
    )


def composition(final_render_path: str | None = "renders/final/alpha.png") -> Composition:
    return Composition(
        composition_id="alpha__20260525",
        target_id="alpha",
        capture_date=date(2026, 5, 25),
        view=ViewState(center=[106.7, 10.8], scale=50000),
        reviewed=True,
        ready=True,
        include=True,
        needs_revalidation=False,
        review_order=1,
        artifacts=CompositionArtifacts(final_render_path=final_render_path),
        layers=[
            ImageLayer(
                layer_id="l1",
                source_path="l1.tif",
                order=0,
                visible=True,
                capture_date=date(2026, 5, 25),
                capture_time=time(8, 30),
                metadata_status=MetadataStatus.VALID,
            )
        ],
    )


def workspace(
    tmp_path: Path,
    final_render_path: str | None = "renders/final/alpha.png",
) -> WorkspaceService:
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    render_dir = service.paths.root / "renders" / "final"
    render_dir.mkdir(parents=True, exist_ok=True)
    if final_render_path is not None:
        (render_dir / "alpha.png").write_bytes(b"png")
    service.write_composition(composition(final_render_path))
    return service


def success_render(spec: RenderSpec, is_cancelled=None) -> RasterRenderResult:
    return RasterRenderResult(
        canvas=np.zeros((spec.output_height, spec.output_width, 3), dtype=np.uint8),
        painted_layer_ids=tuple(layer.layer_id for layer in spec.visible_layers),
    )


def test_export_mode_runs_preflight_and_enables_export_when_ready(tmp_path: Path) -> None:
    qapp()
    mode = ExportMode()
    service = workspace(tmp_path, final_render_path=None)
    target = target_config()
    ensure_final_renders_for_export(service, [target], render=success_render)
    mode.load_workspace(service, targets=[target])

    mode.preflight_button.click()

    assert mode.plan_model.rowCount() == 1
    assert mode.summary.state_label.text() == "Preflight: ready"
    assert mode.export_button.isEnabled() is True
    index = mode.plan_model.index(0, 0)
    assert index.data(ExportPlanRole.COMPOSITION_ID) == "alpha__20260525"
    assert mode.plan_model.index(0, 4).data(Qt.ItemDataRole.DisplayRole) == "0 issues"


def test_export_mode_preflight_includes_template_issues_from_config(tmp_path: Path) -> None:
    qapp()
    mode = ExportMode()
    service = workspace(tmp_path, final_render_path=None)
    target = target_config()
    ensure_final_renders_for_export(service, [target], render=success_render)
    warning = Issue(
        issue_id="target.template_compatibility_unknown",
        severity=IssueSeverity.WARNING,
        scope=IssueScope.TEMPLATE,
        message="Nhieu PPTX template co base/theme/master khac nhau.",
        remediation="Kiem tra base/theme/master truoc khi export.",
    )
    mode.load_workspace(service, targets=[target], template_issues=[warning])

    mode.preflight_button.click()

    assert mode.summary.state_label.text() == "Preflight: warning"
    assert mode.export_button.isEnabled() is True
    assert mode.plan_model.index(0, 4).data(Qt.ItemDataRole.DisplayRole) == "1 issue"
    assert "target.template_compatibility_unknown" in {
        issue.issue_id for issue in mode._last_plan.issues
    }


def test_export_mode_blocks_export_and_exposes_jump_signal(tmp_path: Path) -> None:
    qapp()
    mode = ExportMode(
        export_runner=lambda workspace_service, targets, **kwargs: run_full_export(
            workspace_service,
            targets,
            render=success_render,
            **kwargs,
        )
    )
    mode.load_workspace(
        workspace(tmp_path, final_render_path=None),
        targets=[target_config(txt_template="{unknown}")],
    )
    jumps: list[tuple[str, str, str]] = []
    mode.jumpRequested.connect(lambda target, comp, layer: jumps.append((target, comp, layer)))

    mode.preflight_button.click()
    mode._jump_from_index(mode.plan_model.index(0, 0))

    assert mode.summary.state_label.text() == "Preflight: blocked"
    assert "blocking" in mode.export_button.toolTip()
    assert jumps == [("alpha", "alpha__20260525", "")]


def test_export_mode_runs_full_export_pipeline(tmp_path: Path) -> None:
    app = qapp()
    map_id, _text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = target_config_for_template(tmp_path / "templates" / "alpha.pptx", map_id)
    service = workspace(tmp_path, final_render_path=None)
    mode = ExportMode(
        export_runner=lambda workspace_service, targets, **kwargs: run_full_export(
            workspace_service,
            targets,
            render=success_render,
            **kwargs,
        )
    )
    mode.load_workspace(service, targets=[target])

    mode.preflight_button.click()
    mode.export_button.click()
    _wait_until(
        app,
        lambda: "Export xong" in mode.status_label.text() and mode._export_thread is None,
    )

    assert mode.export_button.isEnabled() is True
    assert (service.paths.exports / "report.pptx").is_file()
    assert (service.paths.exports / "report.txt").read_text("utf-8").strip() == "1|alpha|08:30:00"
    assert (service.paths.exports / "report.export-log.json").is_file()
    assert mode.progress_bar.value() == 100


def test_export_mode_displays_worker_progress_updates(tmp_path: Path) -> None:
    app = qapp()
    map_id, _text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = target_config_for_template(tmp_path / "templates" / "alpha.pptx", map_id)
    service = workspace(tmp_path, final_render_path=None)

    def progress_runner(workspace_service, targets, **kwargs):  # noqa: ANN001
        on_progress = kwargs.get("on_progress")
        if on_progress is not None:
            on_progress(
                ExportProgress(
                    stage="history",
                    completed=65,
                    message="Đang ghi DB history (mục tiêu alpha, đã ghi 3/10 mục tiêu).",
                    current=3,
                    item_total=10,
                    target_id="alpha",
                    composition_id="alpha__20260525",
                )
            )
        time_module.sleep(0.05)
        return run_full_export(
            workspace_service,
            targets,
            render=success_render,
            **kwargs,
        )

    mode = ExportMode(export_runner=progress_runner)
    mode.load_workspace(service, targets=[target])

    mode.preflight_button.click()
    mode.export_button.click()
    _wait_until(app, lambda: "Đang ghi DB history" in mode.status_label.text())

    assert mode.progress_bar.value() == 65

    _wait_until(
        app,
        lambda: "Export xong" in mode.status_label.text() and mode._export_thread is None,
    )
    assert mode.progress_bar.value() == 100


def test_export_mode_shows_export_error_details(tmp_path: Path) -> None:
    app = qapp()
    service = workspace(tmp_path, final_render_path=None)
    target = target_config()
    issue = Issue(
        issue_id="render.raster.no_overlap",
        severity=IssueSeverity.ERROR,
        scope=IssueScope.RENDER,
        target_id="alpha",
        composition_id="alpha__20260525",
        message="Khong co layer visible nao phu vung ban do can render.",
        remediation="Kiem tra lai tam ban do, scale hoac du lieu raster.",
    )

    def failed_runner(workspace_service, targets, **kwargs):  # noqa: ANN001
        del workspace_service, targets, kwargs
        return FullExportResult(
            initial_preflight_plan=mode._last_plan,
            final_render_result=ExportFinalRenderResult(issues=[issue]),
            history_sync_result=ExportHistorySyncResult(enabled=False),
            preflight_plan=mode._last_plan.model_copy(
                update={
                    "rows": [
                        mode._last_plan.rows[0].model_copy(update={"issues": [issue]})
                    ],
                    "issues": [issue],
                }
            ),
            pptx_result=ExportPptxResult(),
            txt_result=ExportTxtResult(),
            log_result=ExportLogWriteResult(
                ok=False,
                summary=ExportCompletionSummary(
                    state=ExportCompletionState.FAILURE,
                    error_count=1,
                    log_path="exports/report.export-log.json",
                ),
                issues=[issue],
            ),
        )

    mode = ExportMode(export_runner=failed_runner)
    mode.load_workspace(service, targets=[target])
    mode.preflight_button.click()
    mode.export_button.click()
    _wait_until(
        app,
        lambda: "Export loi" in mode.status_label.text() and mode._export_thread is None,
    )

    assert "Nguyen nhan" in mode.status_label.text()
    assert "Cach xu ly" in mode.status_label.text()
    assert "Khong co layer visible" in mode.status_label.text()


def test_export_mode_uses_persisted_output_stem(tmp_path: Path) -> None:
    app = qapp()
    map_id, _text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = target_config_for_template(tmp_path / "templates" / "alpha.pptx", map_id)
    service = workspace(tmp_path, final_render_path=None)
    preferences = PreferencesService(tmp_path / "preferences.json")
    preferences.update_export_output_stem("custom_report")
    mode = ExportMode(
        preferences_service=preferences,
        export_runner=lambda workspace_service, targets, **kwargs: run_full_export(
            workspace_service,
            targets,
            render=success_render,
            **kwargs,
        ),
    )
    mode.load_workspace(service, targets=[target])

    assert mode.output_stem_input.text() == "custom_report"
    mode.preflight_button.click()
    mode.export_button.click()
    _wait_until(
        app,
        lambda: "Export xong" in mode.status_label.text() and mode._export_thread is None,
    )

    assert (service.paths.exports / "custom_report.pptx").is_file()
    assert preferences.preferences.export.output_stem == "custom_report"


def _write_template(path: Path) -> tuple[int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    text_shape = slide.shapes.add_textbox(Inches(5), Inches(0.75), Inches(4), Inches(0.5))
    text_shape.text = "Alpha"
    presentation.save(path)
    return int(map_shape.shape_id), int(text_shape.shape_id)


def target_config_for_template(template_path: Path, map_id: int) -> TargetConfig:
    placeholder = TemplatePlaceholder(
        field="map",
        element_id=map_id,
        kind=PlaceholderType.MAP_IMAGE,
        required=True,
    )
    return TargetConfig(
        id="alpha",
        name="Alpha Target",
        geojson_file="targets/alpha.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(template_path),
            "txt_line_template": "{slide_number}|{target_id}|{time_label}",
            "placeholders": [placeholder.model_dump(mode="json")],
        },
        metadata={
            "template_metadata": TemplateMetadata(
                template_pptx=str(template_path),
                slide_index=0,
                map_frame=MapFrame(x=36, y=54, width=288, height=216),
                placeholders=[placeholder],
            ).model_dump(mode="json")
        },
    )


def _wait_until(app: QApplication, predicate, *, timeout: float = 5.0) -> None:  # noqa: ANN001
    deadline = time_module.monotonic() + timeout
    while time_module.monotonic() < deadline:
        app.processEvents()
        if predicate():
            return
    msg = "Timed out waiting for Qt condition"
    raise AssertionError(msg)


def test_app_shell_exposes_export_mode_and_jump_switches_to_review(tmp_path: Path) -> None:
    qapp()
    shell = AppShell(preferences_service=PreferencesService(tmp_path / "preferences.json"))

    assert shell.mode_tabs.count() == 5
    assert shell.mode_tabs.tabText(2) == "Export"
    assert shell.mode_tabs.tabText(3) == "Download"
    assert shell.mode_tabs.tabText(4) == "Config"

    shell.mode_tabs.setCurrentWidget(shell.export_mode)
    shell._jump_to_review_context("alpha", "", "")

    assert shell.mode_tabs.currentWidget() is shell.review_edit_mode


def test_app_shell_config_saved_refreshes_loaded_review_and_export_targets(
    tmp_path: Path,
    monkeypatch,
) -> None:
    qapp()
    shell = AppShell(preferences_service=PreferencesService(tmp_path / "preferences.json"))
    workspace = WorkspaceService(tmp_path / "workspace")
    workspace.initialize(config_path="config.json")
    old_target = target_config()
    new_target = target_config().model_copy(update={"name": "Alpha Reloaded"})
    shell.review_edit_mode.load_workspace(workspace, targets=[old_target])
    shell.export_mode.load_workspace(workspace, targets=[old_target])
    config_path = tmp_path / "config.json"

    def fake_load_project_config(path: Path) -> ConfigLoadResult:
        return ConfigLoadResult(
            config_path=path,
            config=ProjectConfig(targets=[new_target]),
            enabled_targets=[new_target],
        )

    monkeypatch.setattr(
        "thucthengay.editor.app_shell.load_project_config",
        fake_load_project_config,
    )

    shell._config_saved(config_path)

    assert shell.setup_mode.config_row.path_field.full_text == str(config_path.resolve())
    assert shell.review_edit_mode._targets == [new_target]
    assert shell.export_mode._targets == [new_target]
    assert shell.export_mode._last_plan is None
    assert "reload target list" in shell.config_mode.downstream_label.text()
