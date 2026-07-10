from __future__ import annotations

import sqlite3
from datetime import date, time
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from thucthengay.export import ExportProgress, run_full_export
from thucthengay.history import HistoryService
from thucthengay.models import (
    Composition,
    GridConfig,
    GridInterval,
    ImageLayer,
    Issue,
    IssueScope,
    IssueSeverity,
    MapFrame,
    MetadataStatus,
    PlaceholderType,
    TargetConfig,
    TemplateMetadata,
    TemplatePlaceholder,
    TemporalCompareOrientation,
    TemporalCompareState,
    ViewState,
)
from thucthengay.render import RasterRenderResult, RenderError, RenderSpec
from thucthengay.workspace import WorkspaceService


def test_run_full_export_renders_final_images_and_writes_outputs(tmp_path: Path) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())

    result = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
    )

    updated = service.read_composition("alpha__20260525")
    assert result.ok is True
    assert updated.artifacts.final_render_path is not None
    assert updated.artifacts.final_render_path.endswith(".jpg")
    assert (service.paths.root / updated.artifacts.final_render_path).is_file()
    assert (service.paths.exports / "report.pptx").is_file()
    assert (service.paths.exports / "report.txt").read_text("utf-8").strip() == (
        "1|alpha|2026-05-25|08:30:00"
    )
    assert (service.paths.exports / "report.export-log.json").is_file()
    assert result.preflight_plan.summary.error_count == 0


def test_run_full_export_emits_detailed_progress(tmp_path: Path) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())
    history = HistoryService(tmp_path / "history" / "target-history.sqlite")
    events: list[ExportProgress] = []

    result = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        history_service=history,
        on_progress=events.append,
    )

    assert result.ok is True
    assert events[0].stage == "preflight"
    assert events[-1].stage == "done"
    assert events[-1].percent == 100
    messages = "\n".join(event.message for event in events)
    assert "Đang render final image" in messages
    assert "Đang ghi DB history" in messages
    assert "Đang tạo file PPTX" in messages
    assert "Đang tạo file TXT" in messages
    assert "Xong quá trình export" in messages
    assert all(0 <= event.percent <= 100 for event in events)


def test_run_full_export_records_single_pane_history_after_final_preflight(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())
    history = HistoryService(tmp_path / "history" / "target-history.sqlite")

    result = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        history_service=history,
    )

    assert result.ok is True
    assert result.history_sync_result.recorded_layers == 1
    assert result.history_sync_result.include_events == 1
    with sqlite3.connect(history.database_path) as connection:
        link_rows = connection.execute(
            "SELECT latest_composition_id, active FROM target_image_history"
        ).fetchall()
        event_rows = connection.execute("SELECT composition_id FROM include_event").fetchall()
    assert link_rows == [("alpha__20260525", 1)]
    assert event_rows == [("alpha__20260525",)]


def test_run_full_export_copies_sources_to_managed_root_before_history_sync(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    managed_root = tmp_path / "managed-sources"
    target = target.model_copy(
        update={
            "export": target.export.model_copy(
                update={"managed_source_root": str(managed_root)}
            )
        }
    )
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    (service.paths.root / "alpha.tif").write_bytes(b"exported source")
    service.write_composition(_composition())
    history = HistoryService(tmp_path / "history" / "target-history.sqlite")

    result = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        history_service=history,
    )

    assert result.ok is True
    managed_files = list((managed_root / "alpha" / "20260525").glob("alpha__*.tif"))
    assert len(managed_files) == 1
    assert managed_files[0].read_bytes() == b"exported source"
    with sqlite3.connect(history.database_path) as connection:
        rows = connection.execute(
            "SELECT source_path, cache_path, capture_date, capture_time FROM image_asset"
        ).fetchall()
    assert rows == [
        (str(managed_files[0]), None, "2026-05-25", "08:30:00"),
    ]


def test_run_full_export_reuses_existing_managed_source_without_recoping(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    managed_root = tmp_path / "managed-sources"
    target = target.model_copy(
        update={
            "export": target.export.model_copy(
                update={"managed_source_root": str(managed_root)}
            )
        }
    )
    managed_source = managed_root / "alpha" / "20260525" / "alpha__stable.tif"
    managed_source.parent.mkdir(parents=True)
    managed_source.write_bytes(b"already managed source")
    composition = _composition()
    composition = composition.model_copy(
        update={
            "layers": [
                composition.layers[0].model_copy(
                    update={"source_path": str(managed_source)}
                )
            ]
        }
    )
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(composition)
    history = HistoryService(tmp_path / "history" / "target-history.sqlite")

    result = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        history_service=history,
    )

    managed_files = list((managed_root / "alpha" / "20260525").glob("*.tif"))
    assert result.ok is True
    assert managed_files == [managed_source]
    with sqlite3.connect(history.database_path) as connection:
        rows = connection.execute("SELECT source_path FROM image_asset").fetchall()
    assert rows == [(str(managed_source),)]


def test_run_full_export_history_sync_is_idempotent_for_repeated_export(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())
    history = HistoryService(tmp_path / "history" / "target-history.sqlite")

    first = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        history_service=history,
    )
    second = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        history_service=history,
    )

    assert first.history_sync_result.include_events == 1
    assert second.history_sync_result.include_events == 0
    assert second.history_sync_result.existing_layers == 1
    with sqlite3.connect(history.database_path) as connection:
        event_count = connection.execute("SELECT COUNT(*) FROM include_event").fetchone()[0]
    assert event_count == 1


def test_run_full_export_records_compare_pane_compositions_in_history(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    pane_a = _composition(
        composition_id="alpha__20260525",
        capture_date=date(2026, 5, 25),
        include=False,
        review_order=None,
    )
    pane_b = _composition(
        composition_id="alpha__20260526",
        capture_date=date(2026, 5, 26),
        include=False,
        review_order=None,
    )
    selected = _composition(
        composition_id="alpha__20260527",
        capture_date=date(2026, 5, 27),
        review_order=1,
        temporal_compare=TemporalCompareState(
            enabled=True,
            orientation=TemporalCompareOrientation.VERTICAL,
            pane_a_composition_id=pane_a.composition_id,
            pane_b_composition_id=pane_b.composition_id,
        ),
    )
    service.write_composition(pane_a)
    service.write_composition(pane_b)
    service.write_composition(selected)
    history = HistoryService(tmp_path / "history" / "target-history.sqlite")

    result = run_full_export(
        service,
        [target],
        output_stem="compare",
        render=_success_render,
        history_service=history,
    )

    assert result.ok is True
    assert result.history_sync_result.include_events == 2
    with sqlite3.connect(history.database_path) as connection:
        event_rows = connection.execute(
            "SELECT composition_id FROM include_event ORDER BY composition_id"
        ).fetchall()
    assert event_rows == [("alpha__20260525",), ("alpha__20260526",)]


def test_run_full_export_sanitizes_output_stem_for_windows(tmp_path: Path) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())

    result = run_full_export(
        service,
        [target],
        output_stem='daily:report*bad',
        render=_success_render,
    )

    assert result.ok is True
    assert (service.paths.exports / "daily_report_bad.pptx").is_file()
    assert (service.paths.exports / "daily_report_bad.txt").is_file()


def test_run_full_export_carries_template_issues_into_preflight_and_pptx_result(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())
    warning = Issue(
        issue_id="target.template_compatibility_unknown",
        severity=IssueSeverity.WARNING,
        scope=IssueScope.TEMPLATE,
        message="Nhieu PPTX template co base/theme/master khac nhau.",
        remediation="Kiem tra base/theme/master truoc khi export.",
    )

    result = run_full_export(
        service,
        [target],
        output_stem="report",
        render=_success_render,
        template_issues=[warning],
    )

    assert result.ok is True
    assert "target.template_compatibility_unknown" in {
        issue.issue_id for issue in result.preflight_plan.issues
    }
    assert "target.template_compatibility_unknown" in {
        issue.issue_id for issue in result.pptx_result.issues
    }


def test_run_full_export_skips_failed_final_render_without_cascading(
    tmp_path: Path,
) -> None:
    map_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id)
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    service.write_composition(_composition())
    service.write_composition(
        _composition(
            composition_id="alpha__20260526",
            capture_date=date(2026, 5, 26),
            review_order=2,
        )
    )

    result = run_full_export(
        service,
        [target],
        output_stem="partial",
        render=_render_with_no_overlap_failure,
    )

    assert result.ok is True
    assert result.pptx_result.ok is True
    assert result.txt_result.ok is True
    assert [row.composition_id for row in result.pptx_result.exported] == ["alpha__20260525"]
    assert [row.composition_id for row in result.txt_result.exported] == ["alpha__20260525"]
    assert (service.paths.exports / "partial.pptx").is_file()
    assert (service.paths.exports / "partial.txt").read_text("utf-8").strip() == (
        "1|alpha|2026-05-25|08:30:00"
    )
    assert result.log_result is not None
    assert result.log_result.summary.skipped_count == 1
    assert result.log_result.log is not None
    skipped = result.log_result.log.skipped[0]
    assert skipped.composition_id == "alpha__20260526"
    assert "Khong co layer visible" in skipped.reason
    issue_ids = {issue.issue_id for issue in result.log_result.log.issues}
    assert "render.raster.no_overlap" in issue_ids
    assert "export.output_row_missing" not in issue_ids


def _write_template(path: Path) -> int:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    presentation.save(path)
    return int(map_shape.shape_id)


def _target(template_path: Path, map_id: int) -> TargetConfig:
    placeholder = TemplatePlaceholder(
        field="map",
        element_id=map_id,
        kind=PlaceholderType.MAP_IMAGE,
        required=True,
    )
    return TargetConfig(
        id="alpha",
        name="Alpha Target",
        geojson_file="targets/alpha.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(template_path),
            "txt_line_template": "{slide_number}|{target_id}|{capture_date}|{time_label}",
            "placeholders": [placeholder.model_dump(mode="json")],
        },
        metadata={
            "template_metadata": TemplateMetadata(
                template_pptx=str(template_path),
                slide_index=0,
                map_frame=MapFrame(x=36, y=54, width=288, height=216),
                placeholders=[placeholder],
            ).model_dump(mode="json")
        },
    )


def _composition(
    *,
    composition_id: str = "alpha__20260525",
    capture_date: date = date(2026, 5, 25),
    review_order: int | None = 1,
    include: bool = True,
    temporal_compare: TemporalCompareState | None = None,
) -> Composition:
    return Composition(
        composition_id=composition_id,
        target_id="alpha",
        capture_date=capture_date,
        view=ViewState(center=[106.7, 10.8], scale=50000),
        reviewed=True,
        ready=True,
        include=include,
        needs_revalidation=False,
        review_order=review_order,
        temporal_compare=temporal_compare or TemporalCompareState(),
        layers=[
            ImageLayer(
                layer_id=f"{composition_id}-layer",
                source_path="alpha.tif",
                order=0,
                visible=True,
                capture_date=capture_date,
                capture_time=time(8, 30),
                metadata_status=MetadataStatus.VALID,
            )
        ],
    )


def _success_render(spec: RenderSpec, is_cancelled=None) -> RasterRenderResult:
    return RasterRenderResult(
        canvas=np.full((spec.output_height, spec.output_width, 3), 128, dtype=np.uint8),
        painted_layer_ids=tuple(layer.layer_id for layer in spec.visible_layers),
    )


def _render_with_no_overlap_failure(
    spec: RenderSpec,
    is_cancelled=None,  # noqa: ANN001
) -> RasterRenderResult:
    if spec.composition_id == "alpha__20260526":
        raise RenderError(
            [
                Issue(
                    issue_id="render.raster.no_overlap",
                    severity=IssueSeverity.ERROR,
                    scope=IssueScope.RENDER,
                    target_id=spec.target_id,
                    composition_id=spec.composition_id,
                    message="Khong co layer visible nao phu vung ban do can render.",
                    remediation="Kiem tra lai tam ban do, scale hoac du lieu raster.",
                )
            ]
        )
    return _success_render(spec, is_cancelled=is_cancelled)
