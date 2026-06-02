from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from thucthengay.config import load_project_config, update_target_alignment_defaults
from thucthengay.export.final_render import final_render_output_size
from thucthengay.models import GridInterval


def write_json(path: Path, data: dict[str, object]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data), encoding="utf-8")


def write_template_pptx(
    path: Path,
    *,
    slide_count: int = 1,
    width: int | None = None,
    height: int | None = None,
    shape_count: int = 1,
) -> int:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    shape_width = Inches(4) if width is None else width
    shape_height = Inches(3) if height is None else height
    first_shape_id = 0
    for index in range(slide_count):
        slide = presentation.slides.add_slide(blank_layout)
        for shape_index in range(shape_count):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + shape_index),
                Inches(1 + shape_index),
                shape_width,
                shape_height,
            )
            if index == 0 and shape_index == 0:
                first_shape_id = int(shape.shape_id)
    presentation.save(path)
    return first_shape_id


def target_config(
    target_id: str,
    sort_order: int,
    *,
    enabled: bool = True,
    template_pptx_file: str | None = None,
    map_element_id: int = 2,
) -> dict[str, object]:
    return {
        "id": target_id,
        "enabled": enabled,
        "sort_order": sort_order,
        "name": target_id,
        "geojson_file": f"targets/{target_id}.geojson",
        "coordinate": [106.7, 10.8],
        "scale": 50000,
        "grid": {"interval": {"minutes": 1}},
        "export": {
            "template_pptx_file": template_pptx_file or f"templates/{target_id}.pptx",
            "placeholders": [
                {"field": "map_image", "kind": "map_image", "element_id": map_element_id}
            ],
        },
    }


def prepare_target_files(root: Path, target_id: str, template_name: str | None = None) -> int:
    (root / "targets").mkdir(parents=True, exist_ok=True)
    (root / "targets" / f"{target_id}.geojson").write_text("{}", encoding="utf-8")
    return write_template_pptx(root / "templates" / (template_name or f"{target_id}.pptx"))


def write_picture_template_pptx(path: Path, image_path: Path) -> int:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    shape = slide.shapes.add_picture(
        str(image_path),
        Inches(1),
        Inches(1),
        width=Inches(4),
        height=Inches(3),
    )
    presentation.save(path)
    return int(shape.shape_id)


def write_export_contract_template_pptx(path: Path) -> tuple[int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    text_shape = slide.shapes.add_textbox(Inches(5), Inches(0.75), Inches(4), Inches(0.5))
    presentation.save(path)
    return int(map_shape.shape_id), int(text_shape.shape_id)


def write_named_export_contract_template_pptx(path: Path) -> tuple[int, int, int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    map_shape.name = "ttn:map_image"
    title_shape = slide.shapes.add_textbox(Inches(5), Inches(0.75), Inches(4), Inches(0.5))
    title_shape.name = "ttn:title"
    title_shape.text = "NAME, TIME"
    time_shape = slide.shapes.add_textbox(Inches(5), Inches(1.25), Inches(4), Inches(0.5))
    time_shape.name = "ttn:time"
    time_shape.text = "time"
    comment_shape = slide.shapes.add_textbox(Inches(5), Inches(1.75), Inches(4), Inches(0.5))
    comment_shape.name = "ttn:comment"
    comment_shape.text = "comment"
    presentation.save(path)
    return (
        int(map_shape.shape_id),
        int(title_shape.shape_id),
        int(time_shape.shape_id),
        int(comment_shape.shape_id),
    )


def test_load_project_config_filters_enabled_targets_and_resolves_paths(tmp_path: Path) -> None:
    id_b = prepare_target_files(tmp_path, "target_b")
    id_a = prepare_target_files(tmp_path, "target_a")
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                target_config("target_b", 2, map_element_id=id_b),
                target_config("disabled", 0, enabled=False),
                target_config("target_a", 1, map_element_id=id_a),
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    assert [target.id for target in result.enabled_targets] == ["target_a", "target_b"]
    assert result.target_paths["target_a"].geojson_file == (
        tmp_path / "targets" / "target_a.geojson"
    ).resolve()
    assert result.target_paths["target_a"].template_pptx_file == (
        tmp_path / "templates" / "target_a.pptx"
    ).resolve()
    assert result.template_metadata["target_a"].map_frame.width == 288
    assert result.enabled_targets[0].metadata["template_metadata"]["template_pptx"] == str(
        (tmp_path / "templates" / "target_a.pptx").resolve()
    )


def test_load_project_config_accepts_export_contract_from_real_config_shape(
    tmp_path: Path,
) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    map_id, title_id = write_export_contract_template_pptx(
        tmp_path / "templates" / "target_a.pptx"
    )
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                {
                    "id": "target_a",
                    "enabled": True,
                    "sort_order": 1,
                    "name": "Target A",
                    "title": "Target A",
                    "geojson_file": "targets/target_a.geojson",
                    "coordinate": [106.7, 10.8],
                    "scale": 50000,
                    "grid": {"interval": {"minutes": 1}},
                    "export": {
                        "template_pptx_file": "templates/target_a.pptx",
                        "template_txt_value": "Tai {target_title} luc {time_label}",
                        "date_format": "dd.MM.yy",
                        "time_format": "HH.mm/dd.MM.yy",
                        "map_background_color": "#112233",
                        "placeholders": [
                            {"field": "map_image", "element_id": str(map_id)},
                            {
                                "field": "title",
                                "element_id": str(title_id),
                                "value": "Hien trang {target_title} ngay {capture_date}",
                            },
                        ],
                    },
                }
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    target = result.enabled_targets[0]
    assert target.export.template_txt_value == "Tai {target_title} luc {time_label}"
    assert target.export.date_format == "dd.MM.yy"
    assert target.export.time_format == "HH.mm/dd.MM.yy"
    assert target.export.map_background_color == "#112233"
    assert target.export.placeholders[0].kind == "map_image"
    assert target.export.placeholders[0].element_id == map_id
    assert target.export.placeholders[1].kind == "text"
    assert target.export.placeholders[1].value == (
        "Hien trang {target_title} ngay {capture_date}"
    )
    assert result.template_metadata["target_a"].placeholders[1].value == (
        "Hien trang {target_title} ngay {capture_date}"
    )


def test_load_project_config_repairs_stale_placeholder_ids_from_shape_names(
    tmp_path: Path,
) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    map_id, title_id, time_id, comment_id = write_named_export_contract_template_pptx(
        tmp_path / "templates" / "target_a.pptx"
    )
    target = target_config("target_a", 1, map_element_id=999)
    target["export"]["placeholders"] = [  # type: ignore[index]
        {"field": "map_image", "kind": "map_image", "element_id": 999},
        {
            "field": "title",
            "element_id": 998,
            "value": "Hien trang {target_title} ngay {capture_date}",
        },
        {"field": "time", "element_id": 997},
        {"field": "comment", "element_id": 996, "value": "No comment"},
    ]
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    metadata = result.template_metadata["target_a"]
    assert [placeholder.element_id for placeholder in metadata.placeholders] == [
        map_id,
        title_id,
        time_id,
        comment_id,
    ]
    resolution = metadata.metadata["placeholder_resolution"]
    assert [item["configured_element_id"] for item in resolution] == [999, 998, 997, 996]
    assert [item["resolved_element_id"] for item in resolution] == [
        map_id,
        title_id,
        time_id,
        comment_id,
    ]
    assert {item["method"] for item in resolution} == {"field_name"}


def test_load_project_config_resolves_selector_only_placeholder(
    tmp_path: Path,
) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    map_id, *_ = write_named_export_contract_template_pptx(
        tmp_path / "templates" / "target_a.pptx"
    )
    target = target_config("target_a", 1, map_element_id=1)
    target["export"]["placeholders"] = [  # type: ignore[index]
        {
            "field": "map_image",
            "kind": "map_image",
            "selector": {"name": "ttn:map_image"},
        }
    ]
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    placeholder = result.template_metadata["target_a"].placeholders[0]
    assert placeholder.element_id == map_id
    assert placeholder.selector is not None
    assert placeholder.selector.name == "ttn:map_image"


def test_load_project_config_applies_shared_defaults_with_target_overrides(
    tmp_path: Path,
) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "target_a.pptx")
    target = target_config("target_a", 1, map_element_id=element_id)
    target["grid"] = {
        "interval": {"minutes": 1},
        "style": {"label_color": "#445566"},
    }
    target["export"]["time_format"] = "HH:mm"  # type: ignore[index]
    write_json(
        tmp_path / "config.json",
        {
            "defaults": {
                "grid": {
                    "label_format": "dms_short",
                    "style": {
                        "frame_color": "#112233",
                        "label_color": "#000000",
                        "label_font_size": 18,
                        "tick_length_px": 10,
                    },
                },
                "export": {
                    "date_format": "dd.MM.yy",
                    "time_format": "HH.mm/dd.MM.yy",
                    "map_background_color": "#AABBCC",
                },
            },
            "targets": [target],
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    loaded_target = result.enabled_targets[0]
    assert loaded_target.grid.label_format == "dms_short"
    assert loaded_target.grid.style == {
        "frame_color": "#112233",
        "label_color": "#445566",
        "label_font_size": 18,
        "tick_length_px": 10,
    }
    assert loaded_target.export.date_format == "dd.MM.yy"
    assert loaded_target.export.time_format == "HH:mm"
    assert loaded_target.export.map_background_color == "#AABBCC"


def test_pptx_template_metadata_keeps_map_picture_pixel_size(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    image_path = tmp_path / "templates" / "map-placeholder.png"
    image_path.parent.mkdir(parents=True)
    Image.new("RGB", (3306, 2340), (255, 255, 255)).save(image_path)
    element_id = write_picture_template_pptx(tmp_path / "templates" / "target_a.pptx", image_path)
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=element_id)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    metadata = result.enabled_targets[0].metadata["template_metadata"]["metadata"]
    image = metadata["selected_slide"]["shapes"][0]["picture"]["media"]["image"]
    assert image["width_px"] == 3306
    assert image["height_px"] == 2340
    assert final_render_output_size(result.template_metadata["target_a"]) == (3306, 2340)


def test_disabled_targets_are_not_schema_or_reference_blockers(tmp_path: Path) -> None:
    element_id = prepare_target_files(tmp_path, "target_a")
    disabled_invalid_target = {
        "id": "disabled",
        "enabled": False,
        "sort_order": 0,
        "name": "Disabled",
    }
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                disabled_invalid_target,
                target_config("target_a", 1, map_element_id=element_id),
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    assert [target.id for target in result.enabled_targets] == ["target_a"]
    assert "disabled" not in result.target_paths


def test_invalid_config_returns_vietnamese_field_path_issue(tmp_path: Path) -> None:
    data = target_config("target_a", 1)
    data["coordinate"] = [999, 10]
    write_json(tmp_path / "config.json", {"targets": [data]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert result.issues[0].blocking is True
    assert "`targets.0.coordinate`" in result.issues[0].message
    assert "[lon, lat]" in (result.issues[0].remediation or "")


def test_update_target_alignment_defaults_persists_scale_and_grid_interval(
    tmp_path: Path,
) -> None:
    target_a = target_config("target_a", 1)
    target_a["grid"] = {
        "interval": {"minutes": 1},
        "label_format": "dms_short",
        "style": {"frame_color": "#112233"},
    }
    target_b = target_config("target_b", 2)
    config_path = tmp_path / "config.json"
    write_json(config_path, {"targets": [target_a, target_b]})

    updated = update_target_alignment_defaults(
        config_path,
        target_id="target_a",
        interval=GridInterval(minutes=2, seconds=30),
        scale=25000,
    )

    raw = json.loads(config_path.read_text(encoding="utf-8"))
    raw_a = raw["targets"][0]
    raw_b = raw["targets"][1]
    assert updated.id == "target_a"
    assert updated.scale == 25000
    assert updated.grid.interval.minutes == 2
    assert updated.grid.interval.seconds == 30
    assert raw_a["scale"] == 25000
    assert raw_a["grid"]["interval"] == {"minutes": 2, "seconds": 30}
    assert raw_a["grid"]["label_format"] == "dms_short"
    assert raw_a["grid"]["style"] == {"frame_color": "#112233"}
    assert raw_b["scale"] == 50000
    assert raw_b["grid"]["interval"] == {"minutes": 1}


def test_missing_references_create_blocking_target_issues(tmp_path: Path) -> None:
    write_json(tmp_path / "config.json", {"targets": [target_config("target_a", 1)]})

    result = load_project_config(tmp_path / "config.json")
    issue_ids = {issue.issue_id for issue in result.issues}

    assert result.ok is False
    assert "target.geojson_missing" in issue_ids
    assert "target.template_pptx_missing" in issue_ids
    assert all(issue.blocking for issue in result.issues)


def test_config_directory_path_returns_structured_issue(tmp_path: Path) -> None:
    result = load_project_config(tmp_path)

    assert result.ok is False
    assert result.issues[0].issue_id == "config.file_unreadable"
    assert result.issues[0].blocking is True


def test_template_pptx_resolves_relative_to_config_file(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "nested" / "target_a.pptx")
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                target_config(
                    "target_a",
                    1,
                    template_pptx_file="templates/nested/target_a.pptx",
                    map_element_id=element_id,
                )
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    assert result.target_paths["target_a"].template_pptx_file == (
        tmp_path / "templates" / "nested" / "target_a.pptx"
    ).resolve()
    assert result.template_metadata["target_a"].template_pptx == str(
        (tmp_path / "templates" / "nested" / "target_a.pptx").resolve()
    )


def test_template_pptx_must_have_exactly_one_slide(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "target_a.pptx", slide_count=2)
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=element_id)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {
        "target.template_pptx_slide_count_invalid"
    }


def test_required_element_id_must_exist_in_template(tmp_path: Path) -> None:
    prepare_target_files(tmp_path, "target_a")
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=999)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_element_missing"}


def test_ambiguous_placeholder_selector_is_blocking(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    path = tmp_path / "templates" / "target_a.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(2))
    second = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1), Inches(2), Inches(2))
    first.name = "ttn:map_image"
    second.name = "ttn:map_image"
    presentation.save(path)
    target = target_config("target_a", 1, map_element_id=999)
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_element_ambiguous"}


def test_invalid_map_frame_returns_structured_template_issue(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "target_a.pptx", width=0)
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=element_id)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_map_frame_invalid"}


def test_duplicate_placeholder_element_ids_are_blocking(tmp_path: Path) -> None:
    element_id = prepare_target_files(tmp_path, "target_a")
    target = target_config("target_a", 1, map_element_id=element_id)
    target["export"]["placeholders"].append(  # type: ignore[index, union-attr]
        {"field": "target_title", "kind": "text", "element_id": element_id}
    )
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_element_duplicate"}


def test_multiple_required_map_placeholders_are_blocking(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    first_id = write_template_pptx(
        tmp_path / "templates" / "target_a.pptx",
        shape_count=2,
    )
    second_id = first_id + 1
    target = target_config("target_a", 1, map_element_id=first_id)
    target["export"]["placeholders"].append(  # type: ignore[index, union-attr]
        {"field": "map_image_secondary", "kind": "map_image", "element_id": second_id}
    )
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_map_element_ambiguous"}


def test_multiple_template_files_surface_compatibility_warning(tmp_path: Path) -> None:
    id_a = prepare_target_files(tmp_path, "target_a")
    id_b = prepare_target_files(tmp_path, "target_b")
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                target_config("target_a", 1, map_element_id=id_a),
                target_config("target_b", 2, map_element_id=id_b),
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    issue = next(
        issue
        for issue in result.issues
        if issue.issue_id == "target.template_compatibility_unknown"
    )
    assert "target_a" in (issue.remediation or "")
    assert "target_b" in (issue.remediation or "")
