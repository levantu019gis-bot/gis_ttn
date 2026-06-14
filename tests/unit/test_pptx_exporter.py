from __future__ import annotations

from datetime import date, time
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from thucthengay.export import ensure_final_renders_for_export, export_combined_pptx
from thucthengay.models import (
    Composition,
    GridConfig,
    GridInterval,
    ImageLayer,
    MapFrame,
    MetadataStatus,
    PlaceholderType,
    TargetConfig,
    TemplateMetadata,
    TemplatePlaceholder,
    TemporalCompareState,
    ViewState,
)
from thucthengay.render import RasterRenderResult, RenderSpec
from thucthengay.workspace import WorkspaceService


def _write_template(
    path: Path,
    *,
    title: str = "Template",
    static_marker: str | None = None,
) -> tuple[int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    logo_path = path.with_suffix(".logo.png")
    Image.new("RGB", (16, 16), (200, 40, 40)).save(logo_path)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    text_shape = slide.shapes.add_textbox(
        Inches(5),
        Inches(0.75),
        Inches(4),
        Inches(0.5),
    )
    text_shape.text = title
    text_run = text_shape.text_frame.paragraphs[0].runs[0]
    text_run.font.bold = True
    text_run.font.size = Pt(24)
    text_run.font.color.rgb = RGBColor(20, 90, 160)
    if static_marker is not None:
        marker_shape = slide.shapes.add_textbox(
            Inches(5),
            Inches(1.5),
            Inches(4),
            Inches(0.5),
        )
        marker_shape.text = static_marker
    slide.shapes.add_picture(str(logo_path), Inches(9), Inches(0.25), width=Inches(0.4))
    presentation.save(path)
    return int(map_shape.shape_id), int(text_shape.shape_id)


def _write_contract_template(path: Path) -> tuple[int, int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    title_shape = slide.shapes.add_textbox(Inches(5), Inches(0.75), Inches(4), Inches(0.5))
    title_shape.text = "title"
    time_shape = slide.shapes.add_textbox(Inches(5), Inches(1.25), Inches(4), Inches(0.5))
    time_shape.text = "time"
    presentation.save(path)
    return int(map_shape.shape_id), int(title_shape.shape_id), int(time_shape.shape_id)


def _write_contract_template_with_time_rectangle(path: Path) -> tuple[int, int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    title_shape = slide.shapes.add_textbox(Inches(5), Inches(0.75), Inches(4), Inches(0.5))
    title_shape.text = "title"
    time_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5),
        Inches(1.25),
        Inches(0.7),
        Inches(0.3),
    )
    time_shape.text = "time"
    time_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    time_shape.text_frame.word_wrap = True
    presentation.save(path)
    return int(map_shape.shape_id), int(title_shape.shape_id), int(time_shape.shape_id)


def _target(
    template_path: Path,
    map_element_id: int,
    text_element_id: int,
    *,
    target_id: str = "alpha",
    text_field: str = "target_title",
    text_required: bool = True,
) -> TargetConfig:
    placeholders = [
        TemplatePlaceholder(
            field="map",
            element_id=map_element_id,
            kind=PlaceholderType.MAP_IMAGE,
            required=True,
        ),
        TemplatePlaceholder(
            field=text_field,
            element_id=text_element_id,
            kind=PlaceholderType.TEXT,
            required=text_required,
        ),
    ]
    return TargetConfig(
        id=target_id,
        name=f"{target_id.title()} Name",
        alias=target_id.upper(),
        title=f"{target_id.title()} Title",
        geojson_file=f"targets/{target_id}.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(template_path),
            "txt_line_template": "{slide_number}|{target_id}|{capture_date}|{time_label}",
            "placeholders": [
                item.model_dump(mode="json")
                for item in placeholders
            ],
        },
        metadata={
            "template_metadata": TemplateMetadata(
                template_pptx=str(template_path),
                slide_index=0,
                map_frame=MapFrame(x=36, y=54, width=288, height=216),
                placeholders=placeholders,
            ).model_dump(mode="json")
        },
    )


def _composition(
    composition_id: str,
    *,
    target_id: str = "alpha",
    capture_date: date = date(2026, 5, 25),
    review_order: int = 1,
    include: bool = True,
    capture_time: time = time(8, 30),
    temporal_compare: TemporalCompareState | None = None,
) -> Composition:
    return Composition(
        composition_id=composition_id,
        target_id=target_id,
        capture_date=capture_date,
        view=ViewState(center=[106.7, 10.8], scale=50000),
        reviewed=True,
        ready=True,
        include=include,
        needs_revalidation=False,
        review_order=review_order,
        temporal_compare=temporal_compare or TemporalCompareState(),
        layers=[
            ImageLayer(
                layer_id=f"{composition_id}-layer",
                source_path=f"{composition_id}.tif",
                cache_path=f"cache/{composition_id}.tif",
                order=0,
                visible=True,
                capture_date=capture_date,
                capture_time=capture_time,
                metadata_status=MetadataStatus.VALID,
            )
        ],
    )


def _workspace(tmp_path: Path, *compositions: Composition) -> WorkspaceService:
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    for composition in compositions:
        service.write_composition(composition)
    return service


def _success_render(spec: RenderSpec, is_cancelled=None) -> RasterRenderResult:
    return RasterRenderResult(
        canvas=np.full((spec.output_height, spec.output_width, 3), 90, dtype=np.uint8),
        painted_layer_ids=tuple(layer.layer_id for layer in spec.visible_layers),
    )


def test_export_combined_pptx_vertical_slice_replaces_map_and_text(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    output_path = service.paths.exports / "report.pptx"
    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is True
    assert result.summary.slide_count == 1
    assert result.pptx_path == "exports/report.pptx"
    assert result.exported[0].composition_id == "alpha__20260525"
    presentation = Presentation(str(output_path))
    assert len(presentation.slides) == 1
    slide = presentation.slides[0]
    assert any(shape.shape_type == 13 for shape in slide.shapes)
    assert slide.shapes[0].shape_type == 13
    assert slide.shapes[1].text == "Alpha Title"
    exported_run = slide.shapes[1].text_frame.paragraphs[0].runs[0]
    assert exported_run.font.bold is True
    assert exported_run.font.size == Pt(24)
    assert exported_run.font.color.rgb == RGBColor(20, 90, 160)
    slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
    assert "Alpha Title" in slide_text


def test_export_combined_pptx_preserves_all_caps_template_text(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx", title="TITLE")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    output_path = service.paths.exports / "report.pptx"
    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is True
    presentation = Presentation(str(output_path))
    assert presentation.slides[0].shapes[1].text == "ALPHA TITLE"


def test_export_combined_pptx_uses_config_placeholder_values_and_formats(
    tmp_path: Path,
) -> None:
    map_id, title_id, time_id = _write_contract_template(tmp_path / "templates" / "alpha.pptx")
    raw_placeholders = [
        {"field": "map_image", "element_id": str(map_id)},
        {
            "field": "title",
            "element_id": str(title_id),
            "value": "Hien trang {target_title} ngay {capture_date}",
        },
        {"field": "time", "element_id": str(time_id)},
    ]
    target = TargetConfig(
        id="alpha",
        name="Alpha Name",
        title="Alpha Title",
        geojson_file="targets/alpha.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(tmp_path / "templates" / "alpha.pptx"),
            "template_txt_value": "Tai {target_title} luc {time_label}",
            "date_format": "dd.MM.yy",
            "time_format": "HH.mm/dd.MM.yy",
            "placeholders": raw_placeholders,
        },
    )
    target.metadata["template_metadata"] = TemplateMetadata(
        template_pptx=str(tmp_path / "templates" / "alpha.pptx"),
        slide_index=0,
        map_frame=MapFrame(x=36, y=54, width=288, height=216),
        placeholders=target.export.placeholders,
    ).model_dump(mode="json")
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    output_path = service.paths.exports / "contract.pptx"
    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is True
    presentation = Presentation(str(output_path))
    slide_text = "\n".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )
    assert "Hien trang Alpha Title ngay 25.05.26" in slide_text
    assert "08.30/25.05.26" in slide_text
    assert any(shape.shape_type == 13 for shape in presentation.slides[0].shapes)


def test_export_combined_pptx_uses_compare_pane_time_placeholders(
    tmp_path: Path,
) -> None:
    map_id, title_id, time_id = _write_contract_template(tmp_path / "templates" / "alpha.pptx")
    target = TargetConfig(
        id="alpha",
        name="Alpha Name",
        title="Alpha Title",
        geojson_file="targets/alpha.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(tmp_path / "templates" / "alpha.pptx"),
            "template_txt_value": "Tai {target_title} luc {time_label}",
            "date_format": "dd.MM.yy",
            "time_format": "HH.mm/dd.MM.yy",
            "placeholders": [
                {"field": "map_image", "element_id": str(map_id)},
                {
                    "field": "title",
                    "element_id": str(title_id),
                    "value": (
                        "A {capture_date_A} {time_label_pane_A} | "
                        "B {capture_date_B} {time_label_pane_B}"
                    ),
                },
                {"field": "time", "element_id": str(time_id)},
            ],
        },
    )
    target.metadata["template_metadata"] = TemplateMetadata(
        template_pptx=str(tmp_path / "templates" / "alpha.pptx"),
        slide_index=0,
        map_frame=MapFrame(x=36, y=54, width=288, height=216),
        placeholders=target.export.placeholders,
    ).model_dump(mode="json")
    service = _workspace(
        tmp_path,
        _composition(
            "alpha__20260525",
            temporal_compare=TemporalCompareState(
                enabled=True,
                pane_a_composition_id="alpha__20260524",
                pane_b_composition_id="alpha__20260526",
            ),
        ),
        _composition(
            "alpha__20260524",
            capture_date=date(2026, 5, 24),
            capture_time=time(7, 15),
            include=False,
        ),
        _composition(
            "alpha__20260526",
            capture_date=date(2026, 5, 26),
            capture_time=time(9, 45),
            include=False,
        ),
    )
    ensure_final_renders_for_export(service, [target], render=_success_render)

    output_path = service.paths.exports / "compare-time.pptx"
    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is True
    presentation = Presentation(str(output_path))
    slide_text = "\n".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )
    assert "A 24.05.26 07.15/24.05.26 | B 26.05.26 09.45/26.05.26" in slide_text


def test_export_combined_pptx_time_placeholder_rectangle_fits_replaced_text(
    tmp_path: Path,
) -> None:
    map_id, title_id, time_id = _write_contract_template_with_time_rectangle(
        tmp_path / "templates" / "alpha.pptx"
    )
    target = TargetConfig(
        id="alpha",
        name="Alpha Name",
        title="Alpha Title",
        geojson_file="targets/alpha.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(tmp_path / "templates" / "alpha.pptx"),
            "template_txt_value": "Tai {target_title} luc {time_label}",
            "date_format": "dd.MM.yy",
            "time_format": "HH.mm/dd.MM.yy",
            "placeholders": [
                {"field": "map_image", "element_id": str(map_id)},
                {"field": "title", "element_id": str(title_id)},
                {"field": "time", "element_id": str(time_id)},
            ],
        },
    )
    target.metadata["template_metadata"] = TemplateMetadata(
        template_pptx=str(tmp_path / "templates" / "alpha.pptx"),
        slide_index=0,
        map_frame=MapFrame(x=36, y=54, width=288, height=216),
        placeholders=target.export.placeholders,
    ).model_dump(mode="json")
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    output_path = service.paths.exports / "time-fit.pptx"
    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is True
    presentation = Presentation(str(output_path))
    time_shape = next(
        shape
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and shape.text == "08.30/25.05.26"
    )
    assert time_shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    assert time_shape.text_frame.word_wrap is False


def test_export_combined_pptx_orders_slides_by_review_order(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(
        tmp_path,
        _composition("alpha__20260526", capture_date=date(2026, 5, 26), review_order=2),
        _composition("alpha__20260525", capture_date=date(2026, 5, 25), review_order=1),
    )
    ensure_final_renders_for_export(service, [target], render=_success_render)

    result = export_combined_pptx(
        service,
        [target],
        output_path=service.paths.exports / "ordered.pptx",
    )

    assert result.ok is True
    assert [row.composition_id for row in result.exported] == [
        "alpha__20260525",
        "alpha__20260526",
    ]
    assert [row.slide_number for row in result.exported] == [1, 2]


def test_export_combined_pptx_uses_each_targets_template_when_slide_size_matches(
    tmp_path: Path,
) -> None:
    alpha_map_id, alpha_text_id = _write_template(
        tmp_path / "templates" / "alpha.pptx",
        static_marker="ALPHA STATIC",
    )
    beta_map_id, beta_text_id = _write_template(
        tmp_path / "templates" / "beta.pptx",
        static_marker="BETA STATIC",
    )
    alpha = _target(
        tmp_path / "templates" / "alpha.pptx",
        alpha_map_id,
        alpha_text_id,
        target_id="alpha",
    )
    beta = _target(
        tmp_path / "templates" / "beta.pptx",
        beta_map_id,
        beta_text_id,
        target_id="beta",
    )
    service = _workspace(
        tmp_path,
        _composition("alpha__20260525", target_id="alpha", review_order=1),
        _composition("beta__20260525", target_id="beta", review_order=2),
    )
    ensure_final_renders_for_export(service, [alpha, beta], render=_success_render)
    output_path = service.paths.exports / "mixed-compatible.pptx"

    result = export_combined_pptx(service, [alpha, beta], output_path=output_path)

    assert result.ok is True
    assert [row.target_id for row in result.exported] == ["alpha", "beta"]
    presentation = Presentation(str(output_path))
    assert len(presentation.slides) == 2
    slide_texts = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in presentation.slides
    ]
    assert "ALPHA STATIC" in slide_texts[0]
    assert "BETA STATIC" in slide_texts[1]


def test_export_combined_pptx_blocks_unresolved_required_text_placeholder(
    tmp_path: Path,
) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(
        tmp_path / "templates" / "alpha.pptx",
        map_id,
        text_id,
        text_field="unsupported_required_value",
    )
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    result = export_combined_pptx(
        service,
        [target],
        output_path=service.paths.exports / "blocked.pptx",
    )

    assert result.ok is False
    assert "export.pptx_placeholder_unresolved" in {issue.issue_id for issue in result.issues}
    assert not (service.paths.exports / "blocked.pptx").exists()


def test_export_combined_pptx_blocks_missing_final_png(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))

    result = export_combined_pptx(
        service,
        [target],
        output_path=service.paths.exports / "missing-render.pptx",
    )

    assert result.ok is False
    assert "export.final_render_missing" in {issue.issue_id for issue in result.issues}
    assert not (service.paths.exports / "missing-render.pptx").exists()


def test_export_combined_pptx_rejects_out_of_workspace_output_path(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)
    output_path = tmp_path / "outside.pptx"

    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is False
    assert "export.pptx_write_failed" in {issue.issue_id for issue in result.issues}
    assert not output_path.exists()
